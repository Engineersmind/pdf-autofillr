"""PPTX reader using python-pptx."""
from __future__ import annotations
from uploaddocument.readers.base_reader import DocumentReader


class PptxReader(DocumentReader):
    def read(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            content = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())
            if len(content) > 1:
                slides.append("\n".join(content))
        return "\n\n".join(slides)
